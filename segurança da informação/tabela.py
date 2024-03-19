from pptx import Presentation
from pptx.util import Inches
import os 
from PIL import Image as PILImage

# Diretório onde as imagens das interfaces de aplicativos VPN estão armazenadas
imagens_diretorio = "C:/Users/alisson/OneDrive/Área de Trabalho/UNIT/segurança da informação/imagens"

# Nomes das imagens das interfaces de aplicativos VPN
nomes_imagens_vpn = ["vpn1.png", "vpn2.png", "vpn3.png", "vpn4.png", "vpn5.png"]
table_width = Inches(10)  # Aumentei o tamanho da tabela

# Títulos das tabelas comparativas
titulos_tabelas = ["ExpressVPN", "NordVPN", "CyberGhost", "Surfshark", "Private Internet Access (PIA)"]

# Lista de tópicos de verificação para a tabela comparativa
topicos_verificacao = ["Custos", "Funcionalidade", "Velocidade", "Política de Registros", "Segurança", "Compatibilidade", "Suporte Técnico"]

# Adicionar valores à tabela comparativa para cada VPN
valores_vpn = [
    # ExpressVPN
    ["Alto", "Excelente", "Excelente", "Sem registros", "Excelente", "Windows, Mac, iOS, Android", "24/7"],
    
    # NordVPN
    ["Médio", "Excelente", "Excelente", "Sem registros", "Excelente", "Windows, Mac, iOS, Android", "24/7"],
    
    # CyberGhost
    ["Baixo", "Boa", "Boa", "Alguns registros", "Boa", "Windows, Mac, iOS, Android", "24/7"],
    
    # Surfshark
    ["Baixo", "Excelente", "Boa", "Sem registros", "Boa", "Windows, Mac, iOS, Android", "24/7"],
    
    # Private Internet Access (PIA)
    ["Médio", "Boa", "Boa", "Sem registros", "Boa", "Windows, Mac, iOS, Android", "24/7"]
]

# Crie um objeto Presentation (arquivo PPTX)
presentation = Presentation()

# Capa do relatório
titulo_capa = "Relatório de Comparação de VPNs"

# Slide da capa
slide_capa = presentation.slides.add_slide(presentation.slide_layouts[5])  # Layout da capa
title = slide_capa.shapes.title
subtitle = slide_capa.placeholders[1]
title.text = titulo_capa
subtitle.text = "Universidade: Universidade Tiradentes\nCurso: Análise e Desenvolvimento de Sistemas\nDisciplina: Segurança da Informação\nProfessor: Manoel Macedo\nIntegrantes da Equipe: Alisson A. Nunes, Luiz Gustavo Curvelo, José Wellington Ferreira\nTítulo da Tarefa: Melhores VPNs do Mercado"

# Adicionar imagens das interfaces de aplicativos VPN e tabelas comparativas
for i in range(5):
    # Converter imagem WEBP para JPEG
    imagem_webp_path = os.path.join(imagens_diretorio, nomes_imagens_vpn[i])
    imagem_jpeg_path = os.path.splitext(imagem_webp_path)[0] + ".jpeg"
    PILImage.open(imagem_webp_path).convert("RGB").save(imagem_jpeg_path, "JPEG")

    # Slide com imagem da interface do aplicativo VPN
    slide_imagem = presentation.slides.add_slide(presentation.slide_layouts[5])  # Layout de título e conteúdo
    title = slide_imagem.shapes.title
    title.text = titulos_tabelas[i]
    left = Inches(1)
    top = Inches(1)
    height = Inches(4.5)  # Aumentei a altura da imagem
    picture = slide_imagem.shapes.add_picture(imagem_jpeg_path, left, top, height=height)

    # Dados da tabela comparativa (substitua com seus próprios dados)
    dados_tabela = [["Critério"] + topicos_verificacao,
                    ["ExpressVPN"] + valores_vpn[0],
                    ["NordVPN"] + valores_vpn[1],
                    ["CyberGhost"] + valores_vpn[2],
                    ["Surfshark"] + valores_vpn[3],
                    ["Private Internet Access (PIA)"] + valores_vpn[4]]

    # Slide com tabela comparativa
    slide_tabela = presentation.slides[i + 1]  # Obter o slide correspondente à VPN
    table = slide_tabela.shapes.add_table(rows=len(dados_tabela), cols=len(dados_tabela[0]), left=left, top=top, width=table_width, height=height).table
    table.columns[0].width = Inches(2.5)
    
    for row in range(len(dados_tabela)):
        for col in range(len(dados_tabela[0])):
            cell = table.cell(row, col)
            cell.text = dados_tabela[row][col] if col == 0 else valores_vpn[i][col - 1]
            cell.text_frame.paragraphs[0].font.size = Inches(0.2)  # Aumentei o tamanho da fonte

# Slide com o memorial (substitua com seu próprio memorial)
slide_memorial = presentation.slides.add_slide(presentation.slide_layouts[5])  # Layout de título e conteúdo
title = slide_memorial.shapes.title
title.text = "Memorial Descritivo"
left = top = Inches(1)
width = height = Inches(8)
txBox = slide_memorial.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
texto_memorial = """
Com base na análise das tabelas comparativas e considerando as necessidades de uma empresa média com negócios nacionais e internacionais, concluímos que a NordVPN é a melhor escolha. Ela oferece um equilíbrio entre custos acessíveis e uma ampla gama de funcionalidades. Além disso, apresentou excelente velocidade, uma política sólida de não registros e altos padrões de segurança. A compatibilidade e o suporte técnico também foram fatores positivos na nossa decisão.
"""
p = tf.add_paragraph()
p.text = texto_memorial
p.space_after = Inches(0.2)  # Aumentei o espaço após o parágrafo

# Salve o arquivo PPTX
output_filename = "relatorio_vpn.pptx"
presentation.save(output_filename)

print(f"O documento PPTX '{output_filename}' foi criado com sucesso.")
